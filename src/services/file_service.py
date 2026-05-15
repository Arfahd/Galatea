"""
File operations service for handling document read/write/edit.
Supports PDF, DOCX, TXT, XLSX, and PPTX formats.
"""

import logging
import aiofiles
import json
from pathlib import Path
from typing import Optional, Any
from io import BytesIO
from copy import deepcopy

from docx import Document
from docx.shared import Pt
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.text.paragraph import Paragraph
from docx.table import Table
from pypdf import PdfReader
import pdfplumber
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from openpyxl import Workbook, load_workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt

from ..config import config
from ..templates.pptx_templates import get_pptx_template

logger = logging.getLogger(__name__)


class FileService:
    """Service for file operations (read, write, edit)."""

    def __init__(self):
        config.ensure_directories()

    def get_user_directory(self, user_id: int) -> Path:
        """Get or create user-specific directory with validation."""
        if not isinstance(user_id, int) or user_id <= 0:
            raise FileServiceError(f"Invalid user ID: {user_id}")

        user_dir = config.USER_FILES_DIR / str(user_id)

        try:
            resolved = user_dir.resolve()
            base_resolved = config.USER_FILES_DIR.resolve()
            if not str(resolved).startswith(str(base_resolved)):
                logger.error(f"Path traversal attempt detected for user {user_id}")
                raise FileServiceError("Invalid user directory path")
        except (OSError, ValueError) as e:
            logger.error(f"Path resolution error for user {user_id}: {e}")
            raise FileServiceError("Invalid user directory path")

        user_dir.mkdir(parents=True, exist_ok=True)
        return user_dir

    # ==================== General File Operations ====================

    async def read_file(self, file_path: Path) -> str:
        """Read content from a file based on its extension."""
        extension = file_path.suffix.lower()

        if extension == ".pdf":
            return await self._read_pdf(file_path)
        elif extension in {".docx", ".doc"}:
            return await self._read_docx(file_path)
        elif extension == ".txt":
            return await self._read_txt(file_path)
        elif extension == ".xlsx":
            return await self._read_xlsx(file_path)
        elif extension == ".pptx":
            return await self._read_pptx(file_path)
        else:
            raise FileServiceError(f"Unsupported file type: {extension}")

    async def read_file_from_bytes(self, file_bytes: bytes, filename: str) -> str:
        """Read content from file bytes."""
        extension = Path(filename).suffix.lower()

        if extension == ".pdf":
            return self._read_pdf_bytes(file_bytes)
        elif extension in {".docx", ".doc"}:
            return self._read_docx_bytes(file_bytes)
        elif extension == ".txt":
            return file_bytes.decode("utf-8", errors="replace")
        elif extension == ".xlsx":
            return self._read_xlsx_bytes(file_bytes)
        elif extension == ".pptx":
            return self._read_pptx_bytes(file_bytes)
        else:
            raise FileServiceError(f"Unsupported file type: {extension}")

    async def write_file(
        self, content: str, filename: str, user_id: int, file_format: str = "txt",
        original_file_path: Optional[Path] = None
    ) -> Path:
        """Write content to a file."""
        user_dir = self.get_user_directory(user_id)
        safe_filename = self._sanitize_filename(filename)

        if file_format == "txt":
            return await self._write_txt(content, safe_filename, user_dir)
        elif file_format == "docx":
            return await self._write_docx(content, safe_filename, user_dir, original_file_path)
        elif file_format == "pdf":
            return await self._write_pdf(content, safe_filename, user_dir)
        elif file_format == "xlsx":
            return await self._write_xlsx(content, safe_filename, user_dir)
        elif file_format == "pptx":
            return await self._write_pptx(content, safe_filename, user_dir)
        else:
            raise FileServiceError(f"Unsupported output format: {file_format}")

    async def save_uploaded_file(
        self, file_bytes: bytes, filename: str, user_id: int
    ) -> Path:
        """Save an uploaded file to user's directory."""
        user_dir = self.get_user_directory(user_id)
        safe_filename = self._sanitize_filename(filename)
        file_path = user_dir / safe_filename

        counter = 1
        while file_path.exists():
            stem = Path(safe_filename).stem
            suffix = Path(safe_filename).suffix
            file_path = user_dir / f"{stem}_{counter}{suffix}"
            counter += 1

        async with aiofiles.open(file_path, "wb") as f:
            await f.write(file_bytes)

        logger.info(f"Saved file: {file_path}")
        return file_path

    async def cleanup_user_directory(self, user_id: int) -> int:
        """Delete all files in user's directory."""
        user_dir = self.get_user_directory(user_id)
        count = 0
        try:
            for file_path in user_dir.iterdir():
                if file_path.is_file():
                    file_path.unlink()
                    count += 1
            logger.info(f"Cleaned up {count} files for user {user_id}")
        except Exception as e:
            logger.error(f"Error cleaning up user directory {user_id}: {e}")
        return count

    def get_file_size_str(self, file_path: Path) -> str:
        """Get human-readable file size."""
        if not file_path.exists():
            return "Unknown"
        size = file_path.stat().st_size
        for unit in ["B", "KB", "MB", "GB"]:
            if size < 1024:
                return f"{size:.1f} {unit}"
            size /= 1024
        return f"{size:.1f} TB"

    # ==================== PDF Operations ====================

    async def _read_pdf(self, file_path: Path) -> str:
        async with aiofiles.open(file_path, "rb") as f:
            content = await f.read()
        return self._read_pdf_bytes(content)

    def _read_pdf_bytes(self, file_bytes: bytes) -> str:
        text_parts = []
        try:
            with pdfplumber.open(BytesIO(file_bytes)) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text() or ""
                    if text.strip():
                        text_parts.append(f"--- Page {page_num} ---\n{text.strip()}")
        except Exception:
            pass
        return "\n\n".join(text_parts) if text_parts else "(No text content found in PDF)"

    async def _write_pdf(self, content: str, filename: str, user_dir: Path) -> Path:
        if not filename.endswith(".pdf"):
            filename += ".pdf"
        file_path = user_dir / filename
        c = canvas.Canvas(str(file_path), pagesize=letter)
        width, height = letter
        c.setFont("Helvetica", 11)
        y_position = height - inch
        lines = content.split("\n")
        for line in lines[:50]: # Limit for basic version
            c.drawString(inch, y_position, line[:90])
            y_position -= 14
            if y_position < inch: break
        c.save()
        return file_path

    # ==================== DOCX Operations ====================

    async def _read_docx(self, file_path: Path) -> str:
        """Read DOCX file."""
        async with aiofiles.open(file_path, "rb") as f:
            content = await f.read()
        return self._read_docx_bytes(content)

    def _read_docx_bytes(self, file_bytes: bytes) -> str:
        """Read DOCX from bytes with structural markers for tables and images."""
        try:
            doc = Document(BytesIO(file_bytes))
            content_parts = []
            
            table_count = 0
            image_count = 0
            
            for child in doc.element.body:
                if isinstance(child, CT_P):
                    para = Paragraph(child, doc)
                    # Add paragraph text
                    content_parts.append(para.text)
                    
                    # Check for images/graphics in paragraph
                    if 'graphic' in para._element.xml:
                        # Add image marker on a NEW line to ensure visibility to AI
                        content_parts.append(f"<<< IMAGE_{image_count} >>>")
                        image_count += 1
                elif isinstance(child, CT_Tbl):
                    table = Table(child, doc)
                    table_text = self._format_table_for_ai(table)
                    # Add table marker and its content
                    content_parts.append(f"<<< TABLE_{table_count} >>>\n{table_text}")
                    table_count += 1
            
            return (
                "\n\n".join(content_parts)
                if content_parts
                else "(No text content found in document)"
            )
        except Exception as e:
            logger.error(f"Error reading DOCX: {e}")
            raise FileServiceError(f"Failed to read DOCX: {str(e)}")

    def _format_table_for_ai(self, table: Table) -> str:
        """Format a table into a text representation for AI context."""
        rows_text = []
        try:
            # Robust iteration to handle merged cells
            for row in table._element.xpath('.//w:tr'):
                cells = []
                for cell in row.xpath('.//w:tc'):
                    cell_para_texts = [p.text for p in cell.xpath('.//w:p') if p.text]
                    cells.append(" ".join(cell_para_texts).strip())
                rows_text.append("\t|\t".join(cells))
        except Exception as e:
            logger.warning(f"Error formatting table for AI: {e}")
            return "(Table data - see structural marker)"
        return "\n".join(rows_text)

    async def _write_docx(self, content: str, filename: str, user_dir: Path, original_file_path: Optional[Path] = None) -> Path:
        """Write DOCX file with structural preservation."""
        if not filename.endswith(".docx"):
            filename += ".docx"
        file_path = user_dir / filename
        doc = Document()
        original_doc = None
        if original_file_path and original_file_path.exists():
            try:
                original_doc = Document(str(original_file_path))
                logger.info(f"Using original doc for structure: {original_file_path}")
            except Exception as e:
                logger.warning(f"Failed to load original doc for structure: {e}")

        from ..utils.markdown_docx import render_markdown_to_docx
        render_markdown_to_docx(doc, content, original_doc=original_doc)
        doc.save(str(file_path))
        return file_path

    # ==================== TXT Operations ====================

    async def _read_txt(self, file_path: Path) -> str:
        try:
            async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
                return await f.read()
        except UnicodeDecodeError:
            async with aiofiles.open(file_path, "r", encoding="latin-1") as f:
                return await f.read()

    async def _write_txt(self, content: str, filename: str, user_dir: Path) -> Path:
        if not filename.endswith(".txt"): filename += ".txt"
        file_path = user_dir / filename
        async with aiofiles.open(file_path, "w", encoding="utf-8") as f:
            await f.write(content)
        return file_path

    # ==================== XLSX Operations ====================

    async def _read_xlsx(self, file_path: Path) -> str:
        wb = load_workbook(file_path)
        result = []
        for sheet in wb.sheetnames:
            result.append(f"=== Sheet: {sheet} ===")
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                result.append("\t".join([str(c) if c is not None else "" for c in row]))
        return "\n".join(result)

    async def _read_xlsx_bytes(self, file_bytes: bytes) -> str:
        wb = load_workbook(BytesIO(file_bytes))
        result = []
        for sheet in wb.sheetnames:
            result.append(f"=== Sheet: {sheet} ===")
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                result.append("\t".join([str(c) if c is not None else "" for c in row]))
        return "\n".join(result)

    async def _write_xlsx(self, content: str, filename: str, user_dir: Path) -> Path:
        if not filename.endswith(".xlsx"): filename += ".xlsx"
        file_path = user_dir / filename
        wb = Workbook()
        ws = wb.active
        for i, line in enumerate(content.split("\n"), 1):
            for j, val in enumerate(line.split("\t"), 1):
                ws.cell(row=i, column=j, value=val)
        wb.save(str(file_path))
        return file_path

    # ==================== PPTX Operations ====================

    async def _read_pptx(self, file_path: Path) -> str:
        prs = Presentation(file_path)
        result = []
        for i, slide in enumerate(prs.slides, 1):
            result.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"): result.append(shape.text)
        return "\n".join(result)

    async def _read_pptx_bytes(self, file_bytes: bytes) -> str:
        prs = Presentation(BytesIO(file_bytes))
        result = []
        for i, slide in enumerate(prs.slides, 1):
            result.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"): result.append(shape.text)
        return "\n".join(result)

    async def _write_pptx(self, content: str, filename: str, user_dir: Path) -> Path:
        if not filename.endswith(".pptx"): filename += ".pptx"
        file_path = user_dir / filename
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Presentation"
        slide.placeholders[1].text = content[:500]
        prs.save(str(file_path))
        return file_path

    def _sanitize_filename(self, filename: str) -> str:
        filename = Path(filename).name
        invalid_chars = '<>:"/\\|?*'
        for char in invalid_chars:
            filename = filename.replace(char, "_")
        return filename

class FileServiceError(Exception):
    pass
